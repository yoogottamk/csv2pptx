from pptx import Presentation

from config import PATH_DATA, PATH_PPTX
from data import get_data_csv
from utils import duplicate_slide, add_images, replace_text_placeholders

pr = Presentation(pptx=PATH_PPTX)
header, rows = get_data_csv(PATH_DATA)

for row in rows:
    new_slide = duplicate_slide(pr)

    # add all images
    success_img = add_images(new_slide, header, row)
    print(f"Added {success_img} images")

    # replace all text placeholders
    success_txt = replace_text_placeholders(new_slide, header, row)
    print(f"Replaced {success_txt} placeholders")


pr.save("gen.pptx")
print("Done")
